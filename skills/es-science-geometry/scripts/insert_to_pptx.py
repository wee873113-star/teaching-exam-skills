#!/usr/bin/env python3
"""
insert_to_pptx.py — 將幾何圖形（PNG）插入 PowerPoint 投影片

Usage (library - 供其他技能呼叫):
  from insert_to_pptx import insert_figure_to_slide, add_geometry_slide

Usage (CLI):
  python3 insert_to_pptx.py <manifest.json> <target.pptx> [--output out.pptx]
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
import json, sys


def insert_figure_to_slide(prs, slide_idx, png_path,
                            left_cm=3.0, top_cm=3.0, width_cm=10.0):
    """將 PNG 插入指定投影片"""
    slide = prs.slides[slide_idx]
    slide.shapes.add_picture(str(png_path), Cm(left_cm), Cm(top_cm), width=Cm(width_cm))
    return slide


def add_geometry_slide(prs, png_path, title='', subtitle='',
                        left_cm=1.5, top_cm=2.5, width_cm=14.0,
                        layout_idx=6):
    """
    新增一張含幾何圖形的投影片。
    layout_idx=6 為空白版面（大多數主題適用）。
    """
    layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts)-1)]
    slide = prs.slides.add_slide(layout)

    # 標題文字框
    if title:
        txBox = slide.shapes.add_textbox(Cm(1), Cm(0.4), Cm(22), Cm(1.5))
        tf = txBox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True

    # 圖片
    slide.shapes.add_picture(str(png_path), Cm(left_cm), Cm(top_cm), width=Cm(width_cm))

    # 副標題 / 說明
    if subtitle:
        sy = top_cm + width_cm * 0.6 + 0.5   # 圖片下方大約位置
        txBox2 = slide.shapes.add_textbox(Cm(1), Cm(sy), Cm(22), Cm(1.2))
        tf2 = txBox2.text_frame
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(14)
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_figures_comparison_slide(prs, png_paths, labels=None, title='',
                                  top_cm=2.5, width_cm=7.0):
    """
    新增並排比較投影片（最多 3 張圖）。
    """
    layout = prs.slide_layouts[min(6, len(prs.slide_layouts)-1)]
    slide = prs.slides.add_slide(layout)
    n = min(len(png_paths), 3)

    if title:
        txBox = slide.shapes.add_textbox(Cm(1), Cm(0.4), Cm(22), Cm(1.5))
        tf = txBox.text_frame; tf.text = title
        tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.bold = True

    spacing = 25.0 / n
    for i, png in enumerate(png_paths[:n]):
        left = Cm(0.5 + i * spacing)
        slide.shapes.add_picture(str(png), left, Cm(top_cm), width=Cm(width_cm * 0.85))
        if labels and i < len(labels) and labels[i]:
            lx = 0.5 + i*spacing + width_cm*0.85/2 - 1.5
            ly = top_cm + width_cm * 0.55
            txBox = slide.shapes.add_textbox(Cm(lx), Cm(ly), Cm(3), Cm(0.8))
            tf = txBox.text_frame; tf.text = labels[i]
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide


def figures_from_manifest(manifest_path, prs, mode='individual',
                           title_prefix='圖形', width_cm=14.0):
    """
    依 manifest.json 批次插入投影片。
    mode: 'individual'（每圖一頁）或 'comparison'（並排比較頁）
    """
    with open(manifest_path, 'r', encoding='utf-8') as f:
        manifest = json.load(f)
    png_items = [item for item in manifest if 'png' in item]

    if mode == 'comparison':
        pngs    = [item['png']             for item in png_items]
        labels  = [item.get('caption', '') for item in png_items]
        add_figures_comparison_slide(prs, pngs, labels=labels, title=title_prefix, width_cm=8.0)
    else:
        for i, item in enumerate(png_items):
            add_geometry_slide(prs, item['png'],
                               title=item.get('caption', f'{title_prefix} {i+1}'),
                               width_cm=width_cm)


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='插入幾何圖形到 PowerPoint')
    parser.add_argument('manifest', help='manifest.json 路徑')
    parser.add_argument('pptx',     help='目標 .pptx 路徑')
    parser.add_argument('--output', help='輸出路徑')
    parser.add_argument('--mode',   default='individual', choices=['individual','comparison'])
    parser.add_argument('--title',  default='幾何圖形')
    args = parser.parse_args()

    target = Path(args.pptx)
    prs = Presentation(str(target)) if target.exists() else Presentation()
    figures_from_manifest(args.manifest, prs, mode=args.mode, title_prefix=args.title)
    out = args.output or str(target)
    prs.save(out)
    print(f"✅ 已儲存：{out}")
